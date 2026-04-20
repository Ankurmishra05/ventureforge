from io import BytesIO
from fastapi import APIRouter
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

from pptx import Presentation

router = APIRouter()


class ExportRequest(BaseModel):
    idea: str
    startup_name: str
    tagline: str
    market_need: str
    business_model: str
    revenue: int


@router.post("/export/pdf")
def export_pdf(data: ExportRequest):
    buffer = BytesIO()

    pdf = canvas.Canvas(buffer, pagesize=A4)
    pdf.setFont("Helvetica-Bold", 20)
    pdf.drawString(50, 800, data.startup_name)

    pdf.setFont("Helvetica", 12)
    pdf.drawString(50, 770, f"Idea: {data.idea}")
    pdf.drawString(50, 750, f"Tagline: {data.tagline}")
    pdf.drawString(50, 730, f"Market Need: {data.market_need[:80]}")
    pdf.drawString(50, 710, f"Business Model: {data.business_model}")
    pdf.drawString(50, 690, f"Revenue Projection: ${data.revenue}")

    pdf.save()
    buffer.seek(0)

    return StreamingResponse(
        buffer,
        media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=ventureforge_report.pdf"},
    )


@router.post("/export/pitch")
def export_pitch(data: ExportRequest):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = data.startup_name
    slide.placeholders[1].text = data.tagline

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Opportunity"
    slide2.placeholders[1].text = data.market_need

    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Business Model"
    slide3.placeholders[1].text = (
        f"{data.business_model}\nRevenue: ${data.revenue}"
    )

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=ventureforge_pitch.pptx"},
    )